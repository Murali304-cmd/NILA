"""
NILA - RAG (Retrieval-Augmented Generation)
-------------------------------------------
Documents live in the 'documents/' folder and in SQLite; their chunk
embeddings live in a local FAISS index in 'vector_db/'.

Pipeline: upload -> extract text -> clean -> chunk -> embed -> FAISS
          query   -> embed  -> FAISS search -> chunks -> LLM prompt

Everything runs locally. The embedding model is loaded lazily (only when
a document is uploaded or searched) to save RAM on an 8 GB laptop.
"""

import hashlib
import os
import re
import unicodedata

import numpy as np

from . import database

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCUMENTS_DIR = os.path.join(BASE_DIR, "documents")
VECTOR_DIR = os.path.join(BASE_DIR, "vector_db")
INDEX_PATH = os.path.join(VECTOR_DIR, "index.faiss")
IDS_PATH = os.path.join(VECTOR_DIR, "ids.npy")

EMBEDDING_MODEL_NAME = os.environ.get("EMBEDDING_MODEL", "all-MiniLM-L6-v2")

MAX_FILE_SIZE = 20 * 1024 * 1024          # 20 MB
ALLOWED_EXT = {".pdf", ".docx", ".pptx", ".txt", ".md"}
CHUNK_SIZE = 800                          # characters per chunk
CHUNK_OVERLAP = 120

_embedder = None
_index = None
_index_ids = None
_index_sig = None   # (mtime, size) of the index file — reload only if changed


# ---------------------------------------------------------------------------
# Embedding model (lazy, shared)
# ---------------------------------------------------------------------------

def _get_embedder():
    """Load the sentence-transformers model once, on first use."""
    global _embedder
    if _embedder is None:
        try:
            from sentence_transformers import SentenceTransformer
        except ImportError as exc:
            raise RuntimeError(
                "Missing dependency for RAG. Run:  pip install sentence-transformers"
            ) from exc
        _embedder = SentenceTransformer(EMBEDDING_MODEL_NAME)
    return _embedder


def embed_texts(texts):
    """Vectorize a list of strings -> numpy array (normalized for cosine)."""
    vectors = _get_embedder().encode(
        texts, normalize_embeddings=True, show_progress_bar=False,
        batch_size=32, convert_to_numpy=True,
    )
    return np.asarray(vectors, dtype="float32")


def embed_query(text):
    vectors = embed_texts([text])
    return vectors[0]


# ---------------------------------------------------------------------------
# FAISS index
# ---------------------------------------------------------------------------

def _load_index():
    """Load the FAISS index + id list once; reload only if the file changed
    (e.g. rebuilt by a background thread or another process)."""
    global _index, _index_ids, _index_sig
    try:
        sig = (os.path.getmtime(INDEX_PATH), os.path.getsize(INDEX_PATH))
    except OSError:
        if _index is not None:
            _index, _index_ids, _index_sig = None, None, None
        return None
    if _index is not None and sig == _index_sig:
        return _index
    try:
        import faiss
        index = faiss.read_index(INDEX_PATH)
        ids = np.load(IDS_PATH)
        _index, _index_ids, _index_sig = index, ids, sig
    except Exception:
        _index, _index_ids, _index_sig = None, None, None
    return _index


def _save_index(index, ids):
    os.makedirs(VECTOR_DIR, exist_ok=True)
    import faiss
    faiss.write_index(index, INDEX_PATH)
    np.save(IDS_PATH, np.asarray(ids, dtype=np.int64))


def _add_vectors(vectors, ids):
    """Incrementally add new chunk vectors to the FAISS index."""
    global _index, _index_ids, _index_sig
    import faiss
    if vectors.shape[0] == 0:
        return 0
    index = _load_index()
    if index is None:
        index = faiss.IndexFlatIP(vectors.shape[1])
    if index.d != vectors.shape[1]:
        index = faiss.IndexFlatIP(vectors.shape[1])
    index.add(vectors)
    old_ids = list(_index_ids) if _index_ids is not None else []
    all_ids = np.array(old_ids + [int(i) for i in ids], dtype=np.int64)
    _save_index(index, all_ids)
    try:
        _index_sig = (os.path.getmtime(INDEX_PATH), os.path.getsize(INDEX_PATH))
    except OSError:
        _index_sig = None
    _index, _index_ids = index, all_ids
    return vectors.shape[0]


def _build_index():
    """Rebuild FAISS from every chunk currently in the database."""
    global _index, _index_ids, _index_sig
    chunks = database.get_all_chunks()
    if not chunks:
        _index, _index_ids, _index_sig = None, None, None
        if os.path.exists(INDEX_PATH):
            os.remove(INDEX_PATH)
        if os.path.exists(IDS_PATH):
            os.remove(IDS_PATH)
        return 0

    vectors = embed_texts([c["text"] for c in chunks])

    import faiss
    index = faiss.IndexFlatIP(vectors.shape[1])  # cosine, vectors pre-normalized
    index.add(vectors)
    ids = np.array([c["id"] for c in chunks], dtype=np.int64)

    _save_index(index, ids)
    try:
        _index_sig = (os.path.getmtime(INDEX_PATH), os.path.getsize(INDEX_PATH))
    except OSError:
        _index_sig = None
    _index, _index_ids = index, ids
    return len(chunks)


def index_status():
    """Return how many documents/chunks are in the index (for the UI)."""
    docs = database.list_documents()
    chunks = len(database.get_all_chunks())
    return {
        "documents": len(docs),
        "chunks": chunks,
        "index_path": INDEX_PATH,
        "embedding_model": EMBEDDING_MODEL_NAME,
    }


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def _clean_text(text):
    """Strip control chars and collapse repeated blank lines."""
    text = unicodedata.normalize("NFKC", text)
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_text(path, ext):
    """Extract plain text from a supported file."""
    if ext == ".txt" or ext == ".md":
        for enc in ("utf-8", "utf-16", "latin-1"):
            try:
                with open(path, "r", encoding=enc, errors="strict") as f:
                    return _clean_text(f.read())
            except (UnicodeDecodeError, UnicodeError):
                continue
        with open(path, "r", encoding="latin-1") as f:
            return _clean_text(f.read())

    if ext == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise RuntimeError("Missing 'pypdf'. Run:  pip install pypdf") from exc
        reader = PdfReader(path)
        pages = [(page.extract_text() or "") for page in reader.pages]
        return _clean_text("\n\n".join(pages))

    if ext == ".docx":
        try:
            import docx
        except ImportError as exc:
            raise RuntimeError("Missing 'python-docx'. Run:  pip install python-docx") from exc
        doc = docx.Document(path)
        parts = [p.text for p in doc.paragraphs if p.text]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return _clean_text("\n".join(parts))

    if ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("Missing 'python-pptx'. Run:  pip install python-pptx") from exc
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            slide_parts.append(text)
                if shape.shape_type == 19:  # table
                    try:
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells if c.text.strip()]
                            if cells:
                                slide_parts.append(" | ".join(cells))
                    except Exception:
                        pass
            if len(slide_parts) > 1:
                parts.append("\n".join(slide_parts))
        return _clean_text("\n\n".join(parts))

    raise ValueError(f"Unsupported file type: {ext}")


def chunk_text(text):
    """Split text into overlapping chunks for retrieval."""
    text = _clean_text(text)
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))
        # Prefer breaking at a sentence or paragraph boundary.
        if end < len(text):
            boundary = max(text.rfind(". ", start, end),
                           text.rfind("\n", start, end),
                           text.rfind("? ", start, end))
            if boundary > start + CHUNK_SIZE // 2:
                end = boundary + 1
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = max(end - CHUNK_OVERLAP, start + 1)
    return chunks


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def save_file(file_bytes, filename):
    """Validate + persist an uploaded file. Returns the new document row id.
    Fast: no extraction/embedding here (that runs in the background)."""
    ext = os.path.splitext(filename)[1].lower()
    if ext not in ALLOWED_EXT:
        raise ValueError(
            f"Unsupported file type '{ext}'. Allowed: {', '.join(sorted(ALLOWED_EXT))}")
    if len(file_bytes) == 0:
        raise ValueError("The file is empty.")
    if len(file_bytes) > MAX_FILE_SIZE:
        raise ValueError("File is larger than the 20 MB limit.")

    safe_name = os.path.basename(filename).replace("\\", "_").replace("/", "_")
    digest = hashlib.sha1(file_bytes).hexdigest()[:10]
    stored_name = f"{digest}_{safe_name}"
    os.makedirs(DOCUMENTS_DIR, exist_ok=True)
    path = os.path.join(DOCUMENTS_DIR, stored_name)
    with open(path, "wb") as f:
        f.write(file_bytes)

    doc_id = database.add_document(
        filename=safe_name, stored_path=path,
        content_type=ext, size=len(file_bytes))
    return doc_id


def index_document(doc_id):
    """Extract -> chunk -> embed -> add to FAISS. Runs in the background.
    Sets documents.status = ready | failed so the UI can show progress."""
    doc = database.get_document(doc_id)
    if not doc:
        return None
    path = doc["stored_path"]
    ext = os.path.splitext(path)[1].lower()
    try:
        text = extract_text(path, ext)
        chunks = chunk_text(text)
        if not chunks:
            raise ValueError("No readable text was found in the document.")
        database.replace_chunks(doc_id, chunks)
        try:
            vectors = embed_texts(chunks)
            ids = database.get_chunk_ids_for_document(doc_id)
            _add_vectors(vectors, ids)
        except Exception:
            _build_index()  # last resort: full rebuild if incremental add failed
        database.set_document_status(doc_id, "ready")
        return {"id": doc_id, "filename": doc["filename"], "chunks": len(chunks)}
    except Exception as exc:
        database.set_document_status(doc_id, "failed", str(exc)[:300])
        return {"id": doc_id, "filename": doc["filename"], "error": str(exc)[:300]}


def store_document(file_bytes, filename):
    """Synchronous version (validate + extract + index in one call)."""
    doc_id = save_file(file_bytes, filename)
    return index_document(doc_id)


def reindex_document(doc_id):
    """Re-extract a stored file and rebuild its chunks."""
    doc = database.get_document(doc_id)
    if not doc:
        return None
    database.delete_chunks_for_document(doc_id)
    database.set_document_status(doc_id, "processing")
    result = index_document(doc_id)
    _build_index()  # full rebuild after delete to drop old vectors
    return result


def search(query, k=4, document_id=None):
    """Return the k most relevant chunks for a query (empty if no index).
    With document_id set, only chunks of that document are considered."""
    index = _load_index()
    if index is None or index.ntotal == 0:
        return []

    try:
        vector = embed_query(query)
    except Exception:
        return []

    # Filter by document: pull extra candidates, then keep only that doc's.
    probe = max(k * 4, 16) if document_id else k
    probe = min(probe, index.ntotal)
    distances, positions = index.search(vector.reshape(1, -1), probe)

    doc_chunk_ids = set(database.get_chunk_ids_for_document(document_id)) if document_id else None

    results = []
    wanted = []
    for dist, pos in zip(distances[0], positions[0]):
        if pos < 0 or pos >= len(_index_ids):
            continue
        chunk_id = int(_index_ids[pos])
        if doc_chunk_ids is not None and chunk_id not in doc_chunk_ids:
            continue
        wanted.append((chunk_id, dist))
        if len(wanted) >= k:
            break
    chunks_by_id = database.get_chunks_by_ids([c for c, _ in wanted])
    for chunk_id, dist in wanted:
        chunk = chunks_by_id.get(chunk_id)
        if chunk:
            results.append({
                "text": chunk["text"],
                "filename": chunk["filename"],
                "score": round(float(dist), 4),
            })
    return results


def rag_context(query, k=4, document_id=None):
    """Ready-to-inject context string + whether retrieval ran + raw results."""
    results = search(query, k=k, document_id=document_id)
    if not results:
        return "", False, []
    blocks = []
    for i, r in enumerate(results, 1):
        blocks.append(
            f"[Excerpt {i} from '{r['filename']}']\n{r['text']}")
        r["excerpt"] = i
    return "\n\n".join(blocks), True, results


def should_search(message: str) -> bool:
    """Decide if a message needs RAG. Keeps greetings, chat and short
    general-knowledge questions ('What is Python?') OUT of RAG — they never
    need the user's documents and would only add embedding latency."""
    q = message.strip().lower()
    if len(q) < 12:
        return False
    if q in ("hi", "hello", "thanks", "thank you") or q.startswith(("thanks",)):
        return False
    # Strong document words almost always mean RAG.
    if any(w in q for w in ("document", "doc", "pdf", "module", "chapter",
                            "slide", "file", "uploaded", "notes", "notes ",
                            "syllabus", "book", "lecture", "course",
                            "assignment", "summary of the")):
        return True
    # Short questions (<= 40 chars) are general knowledge — skip RAG.
    if len(q) <= 40:
        return False
    # Longer, substantive questions are likely about the user's own material.
    if q.endswith("?") or q.startswith(("explain", "summarize", "summarise",
                                        "describe", "compare", "tell me about",
                                        "what does", "what happened")):
        return True
    return False
